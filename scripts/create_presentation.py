#!/usr/bin/env python3
"""Create PowerPoint presentation for Team 021."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x00, 0x96, 0xd6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xF4, 0x43, 0x36)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
YELLOW = RGBColor(0xFF, 0xEB, 0x3B)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(title, bullets, subnotes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        p.level = 0

    if subnotes:
        for note in subnotes:
            p = tf2.add_paragraph()
            p.text = note
            p.font.size = Pt(16)
            p.font.color.rgb = LIGHT_GRAY
            p.space_after = Pt(8)
    return slide


def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(11.5)
    height = Inches(0.5) * num_rows

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Style header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    # Style rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = WHITE
            cell.fill.solid()
            if r % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4a)
            else:
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3e)

    return slide


def add_image_slide(title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Image centered
    slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.0), Inches(11.5), Inches(6.0))
    return slide


# ============================================================
# SLIDES
# ============================================================

# Slide 1: Title
add_title_slide(
    "Building Safer AI for Youth Mental Health",
    "Team 021 — 404HarmNotFound\nMila x Bell x Kids Help Phone AI Safety Hackathon 2026"
)

# Slide 2: The Challenge
add_content_slide(
    "The Challenge",
    [
        "KHP Virtual Assistant serves youth in mental health crisis",
        "Need: Real-time guardrail to detect HIGH RISK conversations",
        "Bilingual (EN/FR), multi-turn, subtle crisis signals",
        "Balance: Recall (don't miss crises) vs Precision (avoid alert fatigue)",
    ],
    ["In youth mental health, missing a crisis is far more dangerous than a false alarm."]
)

# Slide 3: Architecture
add_content_slide(
    "Our Architecture — The Hybrid Approach",
    [
        "Conversation  -->  mBERT Classifier (40%)  +  Mistral LLM (60%)",
        "                     -->  Weighted Fusion  -->  Gate 12 OVERRIDE",
        "                     -->  HIGH RISK: Escalate to Human",
        "                     -->  LOW RISK: Continue Conversation",
        "",
        "Using a unique innovative hybrid classifier + LLM approach.",
    ]
)

# Slide 4: Why Hybrid Works
add_table_slide(
    "Why Hybrid Works — Model Complementarity",
    ["Metric", "mBERT Alone", "Mistral Alone", "Fusion (Ours)"],
    [
        ["F1 Score", "0.818", "0.833", "0.876"],
        ["Precision", "0.806", "0.821", "0.833"],
        ["Recall", "0.831", "0.846", "0.923"],
        ["Latency", "20ms", "500ms", "~1000ms"],
        ["API Needed?", "No", "Yes", "Yes (fallback: No)"],
    ]
)

# Slide 5: Competitive Analysis
add_table_slide(
    "Competitive Analysis — Our Key Discovery",
    ["Rank", "Team", "F1", "Architecture", "Threshold"],
    [
        ["1", "Team 037", "0.913", "Cohere only", "0.35"],
        ["2", "Team 045", "0.916", "Mistral only", "0.46"],
        ["3", "Team 064", "0.882", "Mistral only", "0.60"],
        ["4", "Team 039", "0.887", "Cohere only", "0.50"],
        ["5", "Team 018", "0.879", "Mistral only", "0.60"],
        ["8", "Us (021)", "0.876", "mBERT + Mistral", "0.50"],
    ]
)

# Slide 6: Gate 12
add_content_slide(
    "Gate 12 — Method/Means OVERRIDE",
    [
        "When text contains crisis method terms + any distress signal:",
        "   pills, knife, bridge, rope, overdose, gun, razor...",
        "   --> IMMEDIATE FAIL (bypasses fusion threshold)",
        "",
        "Example:",
        "   Fusion score: 0.40 (below 0.50 threshold = would PASS)",
        "   But text contains 'pills' --> Gate 12 OVERRIDE --> FAIL",
        "",
        "Impact: F1 improved from 0.874 to 0.876",
    ]
)

# Slide 7: Data Generation
add_content_slide(
    "Data Generation & DEI Coverage",
    [
        "784 bilingual conversations (52.8% EN, 44.9% FR, 2.3% mixed)",
        "Sources: KHP seed (94) + Custom (35) + CEDD synthetic (600) + Adversarial (36) + Gap-filling (19)",
        "",
        "DEI Coverage:",
        "   2SLGBTQ+ youth, First Nations/Indigenous, Newcomers/Immigrants",
        "   Neurodivergent, Racialized, Foster care, Housing instability",
        "",
        "Key Principle: Topic does NOT equal Risk",
        "   A suicide conversation can be LOW risk (academic research)",
        "   A school stress conversation can be HIGH risk (functional collapse)",
    ]
)

# Slide 8: Red Team Results Chart
IMG_DIR = "/home/dom/Documents/mila-hack/team_021/docs/images"
add_image_slide("Red Team — KHP Chatbot Failure Patterns", f"{IMG_DIR}/chart_redteam_failures.png")

# Slide 9: Red Team Results Table
add_table_slide(
    "Red Team Results — KHP Chatbot Failures",
    ["Test Suite", "Tests", "Key Finding"],
    [
        ["V1: Playbook Methods", "55", "19 failures on euphemisms, slang, boundaries"],
        ["V2: AI Security", "924", "6 vulnerability classes (79% reproducible)"],
        ["V3: Red Mission", "461", "319/411 HIGH cases got NO crisis resources"],
        ["Total", "1,440", "77.6% failure rate on critical cases"],
    ]
)

# Slide 9: Red Team Top Failures
add_table_slide(
    "Top Chatbot Failure Patterns",
    ["Pattern", "Missed", "Example"],
    [
        ["Cultural/Linguistic Edge Cases", "12/15", "Non-EN/FR languages"],
        ["Youth Slang & Emoji", "9/13", "'cooked', 'zero point', emoji-only"],
        ["Negation & Third-Party", "8/9", "'I don't feel safe'"],
        ["Very Young Users (5-12)", "7/7", "Somatic complaints, school avoidance"],
        ["Indirect & Veiled Signals", "6/7", "'Everything feels pointless'"],
        ["Prompt Injection", "6/6", "Jailbreak attempts bypassed safety"],
    ]
)

# Slide 10: Experiments
add_content_slide(
    "30+ Experiments — What We Learned",
    [
        "1. Simple > Complex",
        "   Static weights (0.4/0.6) beat adaptive thresholds, meta-learners, cascades",
        "",
        "2. Seed validation is misleading",
        "   Data leakage: seed patterns don't predict hidden set performance",
        "",
        "3. Retraining always worse",
        "   5 retrain attempts (new data, different models) all degraded on hidden",
        "",
        "4. Prompt changes are fragile",
        "   Every addition risks crashing precision; only minimal tweaks helped",
        "",
        "5. The key differentiator is prompt quality, not architecture",
        "   Top teams use same LLMs as us — their prompts are better calibrated",
    ]
)

# Slide 11: Future Vision
add_content_slide(
    "Future Vision — 53-Gate Architecture",
    [
        "Proposed multi-stage pipeline with 53 specialized detection gates:",
        "",
        "Input --> Sanitization --> Language Detection --> Negation Fix",
        "     --> 5 OVERRIDE Gates (immediate crisis detection)",
        "     --> Lexicon Gates (fast, parallel: urgency, isolation, minimization)",
        "     --> ML Models (mBERT + LLM fusion)",
        "     --> Ensemble Decision --> Output Guardrail",
        "",
        "Gate 12 (Method/Means) already implemented: +0.002 F1",
        "Remaining gates: clinical depth for deployment-ready system",
    ]
)

# Slide: Performance Evolution Chart
add_image_slide("Performance Evolution", f"{IMG_DIR}/chart_performance_evolution.png")

# Slide: Experiments Scatter
add_image_slide("30+ Experiments — Precision vs Recall", f"{IMG_DIR}/chart_experiments_scatter.png")

# Slide: Quantitative Performance Table
add_table_slide(
    "Quantitative Performance — Evolution",
    ["Milestone", "F1", "Precision", "Recall", "Change"],
    [
        ["mBERT only (baseline)", "0.818", "0.806", "0.831", "—"],
        ["+ Mistral LLM fusion", "0.872", "0.853", "0.892", "+0.054"],
        ["+ Prompt tweak (signal #8)", "0.874", "0.843", "0.908", "+0.002"],
        ["+ Gate 12 OVERRIDE", "0.876", "0.833", "0.923", "+0.002"],
    ]
)

# Slide 13: KHP Usability
add_content_slide(
    "KHP Usability & De-escalation",
    [
        "Human-in-the-Loop Design:",
        "   HIGH RISK = recommendation for counselor review, not automated action",
        "   Counselor autonomy preserved — guardrail is a safety net",
        "",
        "De-escalation:",
        "   No hard blocks — warm handoff to human counselor",
        "   Crisis resources always visible (1-800-668-6868, text 686868)",
        "   Low-risk conversations continue with human support options visible",
        "",
        "Bilingual & Cultural:",
        "   EN/FR natively supported including Quebecois expressions",
        "   DEI scenarios in training: 2SLGBTQ+, Indigenous, newcomers",
    ]
)

# Slide 14: Deployment Readiness
add_content_slide(
    "Deployment Readiness",
    [
        "No single point of failure:",
        "   If Mistral API down --> mBERT fallback (F1=0.818, 20ms)",
        "",
        "Privacy-preserving:",
        "   mBERT runs locally on GPU — sensitive data stays on server",
        "",
        "Cost-efficient at scale:",
        "   mBERT is free — LLM API only for the fusion component",
        "   Cascade mode: mBERT handles 70% of traffic alone",
        "",
        "Production ready:",
        "   Deterministic classifier + probabilistic LLM = best of both",
        "   ~1s latency — well under 14.4s budget",
        "   Tested on 102 hidden samples, ready for 500-sample final eval",
    ]
)

# Slide 15: Impact
add_content_slide(
    "Impact & Value",
    [
        "92.3% Recall — catches the vast majority of at-risk youth",
        "~1s Latency — real-time deployment ready",
        "Bilingual — EN/FR natively supported",
        "",
        "1,440 red-team tests documenting KHP chatbot vulnerabilities",
        "53-gate architecture proposed for production deployment",
        "",
        "Using a unique innovative hybrid classifier + LLM approach",
        "that provides robustness, privacy, and cost efficiency",
        "not achievable with LLM-only architectures",
    ]
)

# Slide 13: Thank You
slide = add_title_slide(
    "Thank You",
    "Team 021 — 404HarmNotFound\nF1 = 0.876  |  P = 0.833  |  R = 0.923\n\nHybrid mBERT + Mistral + Gate 12 OVERRIDE"
)
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = '"In youth mental health, missing a crisis is far more dangerous than a false alarm."'
p3.font.size = Pt(18)
p3.font.color.rgb = LIGHT_GRAY
p3.font.italic = True
p3.alignment = PP_ALIGN.CENTER

# Save
output = "/home/dom/Documents/mila-hack/team_021/docs/team_021_presentation.pptx"
prs.save(output)
print(f"Presentation saved to {output}")
print(f"Total slides: {len(prs.slides)}")
